"""
Generate the Show-and-Tell PowerPoint deck for Sri Shivasankalpa Vṛnda.
Run: python3 scripts/generate-deck.py
Output: presentation.pptx in the project root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INDIGO = RGBColor(0x2E, 0x2A, 0x5D)
INDIGO_LIGHT = RGBColor(0x5D, 0x57, 0x99)
KUMKUMA = RGBColor(0xA6, 0x32, 0x32)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
CHARCOAL_LIGHT = RGBColor(0x61, 0x61, 0x61)
IVORY = RGBColor(0xF7, 0xF1, 0xE3)
IVORY_DARK = RGBColor(0xED, 0xE5, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color=IVORY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    tf.paragraphs[0].line_spacing = Pt(font_size * line_spacing)
    return tf


def add_para(tf, text, font_size=18, color=CHARCOAL, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri", bullet=False,
             space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(font_size * 0.3)
    if bullet:
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
    return p


def add_divider(slide, top, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), top, Inches(2), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_section_label(slide, text, top=Inches(0.4)):
    add_text_box(slide, Inches(1), top, Inches(4), Inches(0.4),
                 text, font_size=11, color=GOLD, bold=True,
                 font_name="Calibri")


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=CHARCOAL_LIGHT,
                 alignment=PP_ALIGN.RIGHT)


def add_demo_hint(slide, url):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(6.3),
        Inches(4.2), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = INDIGO
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"DEMO  →  open {url} in browser"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


TOTAL = 18

# ════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "तन्मे मनः शिवसङ्कल्पमस्तु",
             font_size=28, color=RGBColor(0xF3, 0xE4, 0xBC), bold=False,
             alignment=PP_ALIGN.CENTER, font_name="Arial Unicode MS")

add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(1.2),
             "Sri Shivasankalpa Vṛnda",
             font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Website v0.1 — Show & Tell",
             font_size=22, color=RGBColor(0xD4, 0xD2, 0xE8), bold=False,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
             "Preserving Sanātana Dharma through Seva, Śraddhā, and Saṁskāra",
             font_size=16, color=RGBColor(0xF3, 0xE4, 0xBC), bold=False,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
             "April 2026", font_size=14, color=RGBColor(0x99, 0x94, 0xBF),
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Why a Website Now?
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 1 — WHY AND HOW")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Why a Website Now?",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4.5),
                  "The Vṛnda's mission deserves a digital presence that matches its dignity.",
                  font_size=20, color=CHARCOAL)

add_para(tf, "", font_size=10)
add_para(tf, "The Maharudra Purascharana (May 15–17) is the catalyst — we need a place where devotees can learn about the event, understand the sevas, and contribute.",
         font_size=18, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "", font_size=10)
add_para(tf, "What the website does:", font_size=18, color=INDIGO, bold=True, space_before=8)
add_para(tf, "    ●  Informs — who we are, what we do, which Gurukulas we support", font_size=17, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Inspires — the importance of Veda Vidya, the Guru–Shishya Parampara", font_size=17, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Invites participation — events, volunteering, Gurukula visits", font_size=17, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Accepts donations — UPI, QR code, bank transfer", font_size=17, color=CHARCOAL_LIGHT, space_before=4)

add_slide_number(slide, 2, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Why This Approach (not WordPress)?
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 1 — WHY AND HOW")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Why Not Just Use WordPress?",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

# WordPress column
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1), Inches(2.3), Inches(5), Inches(4.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF9, 0xEC, 0xEC)
shape.line.fill.background()

tf = add_text_box(slide, Inches(1.3), Inches(2.5), Inches(4.4), Inches(4.0),
                  "Typical WordPress Site", font_size=20, color=KUMKUMA, bold=True)
add_para(tf, "", font_size=6)
add_para(tf, "✗  Generic templates — looks like every other NGO site", font_size=15, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "✗  Plugin bloat — dozens of plugins, each a security risk", font_size=15, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "✗  Needs a database — can be hacked, needs maintenance", font_size=15, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "✗  Monthly hosting cost (₹3,000–10,000/year)", font_size=15, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "✗  Slow on mobile — heavy pages, poor performance", font_size=15, color=CHARCOAL_LIGHT, space_before=6)

# Our approach column
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.2))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0xEE, 0xED, 0xF5)
shape2.line.fill.background()

tf2 = add_text_box(slide, Inches(7.1), Inches(2.5), Inches(5.0), Inches(4.0),
                   "Our Custom-Built Site", font_size=20, color=INDIGO, bold=True)
add_para(tf2, "", font_size=6)
add_para(tf2, "✓  Designed for this trust's aesthetic — not a template", font_size=15, color=CHARCOAL, space_before=6)
add_para(tf2, "✓  No plugins — clean, secure, nothing to maintain", font_size=15, color=CHARCOAL, space_before=6)
add_para(tf2, "✓  No database to hack — static pages, ultra-safe", font_size=15, color=CHARCOAL, space_before=6)
add_para(tf2, "✓  Free hosting on Vercel (₹0/year)", font_size=15, color=CHARCOAL, space_before=6)
add_para(tf2, "✓  Blazing fast on mobile — loads in under 1 second", font_size=15, color=CHARCOAL, space_before=6)

add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
             "\"We built a precision instrument, not a Swiss Army knife.\"",
             font_size=16, color=INDIGO_LIGHT, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Calibri")

add_slide_number(slide, 3, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Design Philosophy
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 1 — WHY AND HOW")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Design Philosophy",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

# Color swatches
colors_data = [
    ("Ivory", IVORY, CHARCOAL), ("Charcoal", CHARCOAL, WHITE),
    ("Indigo", INDIGO, WHITE), ("Kumkuma", KUMKUMA, WHITE),
    ("Gold", GOLD, WHITE),
]
for i, (name, col, txt_col) in enumerate(colors_data):
    x = Inches(1 + i * 2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, Inches(2.2), Inches(1.8), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = col
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.1), Inches(2.45), Inches(1.6), Inches(0.6),
                 name, font_size=14, color=txt_col, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text_box(slide, Inches(1), Inches(3.6), Inches(10), Inches(3.5),
                  "Typography-first:", font_size=18, color=INDIGO, bold=True)
add_para(tf, "    EB Garamond for headings  ·  Inter for body text  ·  Noto Serif Devanagari for shlokas",
         font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "", font_size=10)
add_para(tf, "Core principles:", font_size=18, color=INDIGO, bold=True, space_before=8)
add_para(tf, "    ●  Sanskrit displayed with reverence — Devanagari + transliteration + translation",
         font_size=16, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "    ●  Mobile-first — 60%+ visitors will be on phones",
         font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  No generic stock photos, no kitsch-Hindu design",
         font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Dignified, rooted, contemporary — like a museum exhibit about a living tradition",
         font_size=16, color=CHARCOAL_LIGHT, space_before=4)

add_demo_hint(slide, "localhost:3001/styleguide")
add_slide_number(slide, 4, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Content Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 1 — WHY AND HOW")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "How Content Works",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(1.0),
                  "All content lives in simple text files. No database, no CMS login. Anyone with a text editor can update the site.",
                  font_size=18, color=CHARCOAL_LIGHT)

# File table
files = [
    ("gurukulas.ts", "Gurukula names, locations, descriptions"),
    ("events.ts", "Events — Maharudra schedule, seva amounts"),
    ("trustees.ts", "Trustees & Managing Committee bios"),
    ("blog.ts", "Blog articles"),
    ("gallery.ts", "Gallery photos and videos"),
    ("faqs.ts", "Frequently asked questions"),
]
for i, (fname, desc) in enumerate(files):
    y = Inches(3.4 + i * 0.48)
    add_text_box(slide, Inches(1.5), y, Inches(3), Inches(0.45),
                 fname, font_size=15, color=INDIGO, bold=True, font_name="Courier New")
    add_text_box(slide, Inches(4.8), y, Inches(6), Inches(0.45),
                 desc, font_size=15, color=CHARCOAL_LIGHT)

tf2 = add_text_box(slide, Inches(1), Inches(6.4), Inches(10), Inches(0.5),
                   "Future: we'll migrate to a visual CMS (Sanity) — the site won't change, only the data source.",
                   font_size=15, color=INDIGO_LIGHT)

add_slide_number(slide, 5, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — Homepage
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Homepage",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

sections = [
    "Hero — banner image of Gurukula students, tagline shloka, trust logo, audio player for Vedic chanting",
    "Veda Vidya — Vishnu Purana shloka + the Veda Vriksha illustration showing all 4 Vedas and their shakhas",
    "Why Gurukulas Matter — the Vṛnda's own founding text, with Adi Shankaracharya's injunction",
    "Featured Gurukulas — 3 cards (Shruti Parampara, Gowtama, Sacchidananda) linking to detail pages",
    "Maharudra Purascharana — event card with date, description, and direct link",
    "Blessed by the Sringeri Jagadgurus — anugraha section with the Sringeri Parampara image",
    "Support the Cause — 3 donation purpose cards, each opening a payment modal (live UPI + QR + bank)",
]
for i, s in enumerate(sections):
    parts = s.split(" — ", 1)
    y = Inches(2.2 + i * 0.62)
    add_text_box(slide, Inches(1.3), y, Inches(2.5), Inches(0.55),
                 parts[0], font_size=15, color=INDIGO, bold=True)
    add_text_box(slide, Inches(3.9), y, Inches(8), Inches(0.55),
                 parts[1] if len(parts) > 1 else "", font_size=14, color=CHARCOAL_LIGHT)

add_demo_hint(slide, "localhost:3001/")
add_slide_number(slide, 6, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Gurukulas
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Vedic Gurukulas",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4),
                  "List page — card grid of 3 seeded Gurukulas with name, location, ācārya, student count.",
                  font_size=18, color=CHARCOAL_LIGHT)
add_para(tf, "", font_size=8)
add_para(tf, "Detail pages — each Gurukula has its own URL (/gurukulas/shruti-parampara) designed as a drawer-style page:",
         font_size=18, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "", font_size=6)
add_para(tf, "    ●  Hero image with overlay", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  5 tabs: Overview, Adhyāpakas, Vidyārthīs, Events, Contact", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Close button (×) returns to the list", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  SEO-friendly — each page is linkable and shareable", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "", font_size=8)
add_para(tf, "Seeded Gurukulas:", font_size=18, color=INDIGO, bold=True, space_before=8)
add_para(tf, "    1. Shruti Parampara Gurukula, JP Nagar — 12 students", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    2. Gowtama Veda Pathashala, Banashankari — 20 students", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    3. Sacchidananda Advaitashrama, Rajajinagar — 12 students", font_size=16, color=CHARCOAL_LIGHT, space_before=4)

add_demo_hint(slide, "localhost:3001/gurukulas")
add_slide_number(slide, 7, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — Maharudra
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Maharudra Purascharana",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4.5),
                  "A dedicated page for our flagship event — sourced from the official invite letter.",
                  font_size=18, color=CHARCOAL_LIGHT)
add_para(tf, "", font_size=8)
add_para(tf, "    ●  Sri Adi Shankaracharya painting (Sphatika Linga worship) — full, unclipped", font_size=16, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf, "    ●  Spiritual significance of the Rudrādhyāya with corrected Shiva Purana shlokas", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  3-day programme schedule (May 15–17)", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  10 seva opportunities with amounts (₹501 to ₹15,001)", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Google Maps embed — Hoysala Trust, Hosakerehalli", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  \"Donate\" button opens a payment modal directly (UPI + QR + bank)", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  \"Volunteer\" button links to the contact form", font_size=16, color=CHARCOAL_LIGHT, space_before=4)

add_demo_hint(slide, "localhost:3001/events/maharudra")
add_slide_number(slide, 8, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — Events & Donations
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Events & Donations",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(5), Inches(4),
                  "Events (/events)", font_size=22, color=INDIGO, bold=True)
add_para(tf, "", font_size=6)
add_para(tf, "Introduction to why we organise events, followed by featured event cards.", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "More events will appear here as they are planned.", font_size=16, color=CHARCOAL_LIGHT, space_before=4)

tf2 = add_text_box(slide, Inches(6.5), Inches(2.3), Inches(6), Inches(4),
                   "Donations (/donations)", font_size=22, color=INDIGO, bold=True)
add_para(tf2, "", font_size=6)
add_para(tf2, "\"Śraddhayā deyam\" illustration at top.", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf2, "Three purpose cards:", font_size=16, color=CHARCOAL_LIGHT, space_before=6)
add_para(tf2, "    ●  Gurukula Abhivruddhi", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  Go-Samrakshanam", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  Event Seva", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "", font_size=6)
add_para(tf2, "Each opens a modal with live payment info:", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf2, "UPI ID, scannable QR code, bank details.", font_size=16, color=INDIGO, bold=True, space_before=2)

add_demo_hint(slide, "localhost:3001/donations")
add_slide_number(slide, 9, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — About Us
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "About Us",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4.5),
                  "Our Story — the founding narrative (\"We are not a Gurukulam. We are a trust that exists to ensure Gurukulas thrive.\")",
                  font_size=17, color=CHARCOAL_LIGHT)
add_para(tf, "", font_size=6)
add_para(tf, "Vision & Mission — prominent cards on a deep indigo background", font_size=17, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "What We Do — 6 objectives shown as icon cards (not a bullet list)", font_size=17, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "How We Operate — reaching out to Gurukulas, building a database", font_size=17, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "", font_size=6)
add_para(tf, "Team Shivasankalpa:", font_size=18, color=INDIGO, bold=True, space_before=10)
add_para(tf, "    Managing Committee: Naveen (President), Sheshadri (VP), Shreesha (Secretary), Gourishankara (Treasurer)", font_size=15, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    Trustees: Anantanarayana Sharma, Ravishankar, Giri Bhardwaj, Harisha, Bodhayana", font_size=15, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    Click any card → expanded modal with larger photo and full bio", font_size=15, color=INDIGO_LIGHT, space_before=4)

add_demo_hint(slide, "localhost:3001/about")
add_slide_number(slide, 10, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 11 — Blog
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Blog",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4),
                  "Two seeded posts to demonstrate the layout:", font_size=18, color=CHARCOAL_LIGHT)
add_para(tf, "", font_size=8)
add_para(tf, "  1.  \"A Visit to Shruti Parampara Gurukula, JP Nagar\"", font_size=17, color=INDIGO, bold=True, space_before=8)
add_para(tf, "       A detailed visit report — meeting the ācārya, watching students chant, identifying needs.", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "", font_size=6)
add_para(tf, "  2.  \"Why Vedic Gurukulas Matter in the 21st Century\"", font_size=17, color=INDIGO, bold=True, space_before=8)
add_para(tf, "       A reflective essay on the Gurukula tradition — to be completed by the team.", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "", font_size=10)
add_para(tf, "Each post has: date, author, tags, full-width image, rich text body.", font_size=16, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "The team can add new posts by editing a simple text file (data/blog.ts).", font_size=16, color=CHARCOAL_LIGHT, space_before=4)

add_demo_hint(slide, "localhost:3001/blog")
add_slide_number(slide, 11, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — Gallery
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Gallery",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(4),
                  "31 photos and videos from Gurukula visits, organised with tab filters:", font_size=18, color=CHARCOAL_LIGHT)
add_para(tf, "", font_size=8)
add_para(tf, "    All  ·  Gurukulas  ·  Events  ·  Misc", font_size=18, color=INDIGO, bold=True, space_before=6)
add_para(tf, "", font_size=8)
add_para(tf, "Sources:", font_size=17, color=CHARCOAL_LIGHT, space_before=8)
add_para(tf, "    ●  Shruti Parampara (1 photo, 2 videos)", font_size=15, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "    ●  Varthur Gurukula (12 photos, 4 videos)", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "    ●  Sri Ramana Maharshi Brahmavidyashrama (4 photos, 1 video)", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "    ●  Chidambarashrama (4 photos)", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "    ●  Events (1 photo, 1 video)", font_size=15, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf, "", font_size=8)
add_para(tf, "Videos play inline — click the play button to watch.", font_size=16, color=INDIGO_LIGHT, space_before=6)

add_demo_hint(slide, "localhost:3001/gallery")
add_slide_number(slide, 12, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — Contact & FAQs
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Contact & FAQs",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(5), Inches(4),
                  "Contact Form (/contact)", font_size=22, color=INDIGO, bold=True)
add_para(tf, "", font_size=6)
add_para(tf, "Fields: Name, Email, Subject, Message", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "Client-side validation (required fields, email format)", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "Shows a success confirmation on submit", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf, "Email delivery will be connected at deployment (v0.15)", font_size=15, color=INDIGO_LIGHT, space_before=6)

tf2 = add_text_box(slide, Inches(6.5), Inches(2.3), Inches(6), Inches(4),
                   "FAQs (/faqs)", font_size=22, color=INDIGO, bold=True)
add_para(tf2, "", font_size=6)
add_para(tf2, "8 questions with expandable accordion:", font_size=16, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf2, "    ●  What is Sri Shivasankalpa?", font_size=14, color=CHARCOAL_LIGHT, space_before=4)
add_para(tf2, "    ●  What does \"Shivasankalpa\" mean?", font_size=14, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  How are donations used?", font_size=14, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  How can I participate?", font_size=14, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  Are donations tax-exempt?", font_size=14, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  Which Gurukulas does the trust support?", font_size=14, color=CHARCOAL_LIGHT, space_before=2)
add_para(tf2, "    ●  ...and more", font_size=14, color=CHARCOAL_LIGHT, space_before=2)

add_demo_hint(slide, "localhost:3001/contact")
add_slide_number(slide, 13, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 14 — 404 Page
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 2 — PAGE WALKTHROUGH")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Even the 404 Page Has Character",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(3),
                  "When someone visits a page that doesn't exist, they see:",
                  font_size=18, color=CHARCOAL_LIGHT, alignment=PP_ALIGN.CENTER)
add_para(tf, "", font_size=12)
add_para(tf, "\"404 — Page Not Found\"", font_size=24, color=INDIGO, bold=True,
         alignment=PP_ALIGN.CENTER, space_before=10)
add_para(tf, "", font_size=6)
add_para(tf, "तमेव शरणं गच्छ सर्वभावेन भारत", font_size=20, color=INDIGO_LIGHT,
         alignment=PP_ALIGN.CENTER, font_name="Arial Unicode MS", space_before=10)
add_para(tf, "\"Seek refuge in Him alone with your whole being.\" — Bhagavad Gītā 18.62",
         font_size=15, color=CHARCOAL_LIGHT, alignment=PP_ALIGN.CENTER, space_before=6)
add_para(tf, "", font_size=8)
add_para(tf, "...with a button to return home.", font_size=16, color=CHARCOAL_LIGHT,
         alignment=PP_ALIGN.CENTER, space_before=6)

add_slide_number(slide, 14, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — Placeholders Overview
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 3 — WHAT'S NEXT")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "What Needs Your Input",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.6),
                  "~45 placeholders in the site need real content before we go live. Organised by category:",
                  font_size=18, color=CHARCOAL_LIGHT)

categories = [
    ("Gurukula details", "Established year, history, daily schedule, phone, email — for all 3 Gurukulas"),
    ("Trustee & committee bios", "Final bio paragraphs for all 9 members"),
    ("Team photos", "Formal photos for members who don't have one yet"),
    ("Donation account", "Current details are interim — need trust's own bank/UPI"),
    ("Blog content", "Review the visit post; complete the second article"),
    ("Trust artefacts", "Registration certificate, audit report"),
    ("Trust registration no.", "For the footer"),
    ("80G status", "Update the FAQ once certification is obtained"),
]
for i, (cat, desc) in enumerate(categories):
    y = Inches(3.2 + i * 0.48)
    add_text_box(slide, Inches(1.3), y, Inches(3.5), Inches(0.45),
                 cat, font_size=15, color=INDIGO, bold=True)
    add_text_box(slide, Inches(5.0), y, Inches(7), Inches(0.45),
                 desc, font_size=14, color=CHARCOAL_LIGHT)

add_slide_number(slide, 15, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 16 — Media Review
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 3 — WHAT'S NEXT")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Media Review Needed",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.6),
                  "We need better photos and videos in several areas. Please review and provide upgrades:",
                  font_size=18, color=CHARCOAL_LIGHT)

items = [
    ("Gurukula hero images", "Dedicated photo for each Gurukula's detail page"),
    ("Trustee & committee portraits", "Formal photos, good lighting, consistent style"),
    ("Gallery", "Replace blurry/dark photos; add photos from recent visits"),
    ("Blog images", "Relevant photos for both seeded blog posts"),
    ("Homepage banner", "Consider an original photo from a Vṛnda event"),
    ("Veda Vruksha", "Verify accuracy of shakha names in the illustration"),
]
for i, (area, action) in enumerate(items):
    y = Inches(3.2 + i * 0.55)
    add_text_box(slide, Inches(1.3), y, Inches(3.5), Inches(0.50),
                 area, font_size=15, color=KUMKUMA, bold=True)
    add_text_box(slide, Inches(5.0), y, Inches(7), Inches(0.50),
                 action, font_size=14, color=CHARCOAL_LIGHT)

add_slide_number(slide, 16, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 17 — Roadmap
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ACT 3 — WHAT'S NEXT")

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Roadmap",
             font_size=36, color=INDIGO, bold=True)
add_divider(slide, Inches(1.8))

versions = [
    ("v0.1  (today)", "Prototype on localhost — what you're seeing now", INDIGO),
    ("v0.15", "Deploy to srishivasankalpa.org (Vercel + Hostinger DNS), connect contact form email", INDIGO_LIGHT),
    ("v0.2", "Sanity CMS (visual editor for non-technical updates), Razorpay payment gateway, automated 80G receipts", CHARCOAL),
    ("v1.0", "Multi-language support, calendar-based seva booking, donor dashboard", CHARCOAL_LIGHT),
]
for i, (ver, desc, col) in enumerate(versions):
    y = Inches(2.5 + i * 1.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), y, Inches(11), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xED, 0xF5) if i == 0 else RGBColor(0xFB, 0xF8, 0xF0)
    shape.line.fill.background()

    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(2), Inches(0.7),
                 ver, font_size=18, color=col, bold=True)
    add_text_box(slide, Inches(3.5), y + Inches(0.08), Inches(8), Inches(0.7),
                 desc, font_size=15, color=CHARCOAL_LIGHT)

add_slide_number(slide, 17, TOTAL)

# ════════════════════════════════════════════════════════════
# SLIDE 18 — CTA / Asks
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, INDIGO)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "What We Need From You",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

asks = [
    "Review and provide your final bios (2–3 sentences each)",
    "Send a formal photo for the Team page (good lighting, head-and-shoulders)",
    "Review the 3 Gurukula entries — fill in real details (history, schedule, contact)",
    "Read and approve (or revise) the two blog posts",
    "Provide the trust registration number",
    "Set up the trust's own bank account + UPI — current details are interim",
    "Share feedback on the site's design, tone, and content",
]
for i, ask in enumerate(asks):
    y = Inches(2.0 + i * 0.65)
    num_col = RGBColor(0xF3, 0xE4, 0xBC)
    add_text_box(slide, Inches(2), y, Inches(0.5), Inches(0.55),
                 f"{i+1}.", font_size=18, color=num_col, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(2.7), y, Inches(8), Inches(0.55),
                 ask, font_size=17, color=RGBColor(0xD4, 0xD2, 0xE8))

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "॥ शिवसंकल्पमस्तु ॥",
             font_size=22, color=RGBColor(0xF3, 0xE4, 0xBC), bold=False,
             alignment=PP_ALIGN.CENTER, font_name="Arial Unicode MS")

add_slide_number(slide, 18, TOTAL)

# ════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════
output_path = "/Users/karthik/go/src/github.com/shivasankalpa-stack/website/presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
